"""Academic slide deck: Knowable Worlds motivation + methods.

Generates slides/knowable_worlds_motivation_methods.pptx (16:9) plus the two
figures it embeds (built from the real study worlds — seed 300). Content
covers motivation, related-work positioning, the dynamic world, both question
channels, the two extensions, controls, and registered hypotheses. No
results — this is the motivation/methods deck.

    python -m knowable_worlds.gen_slides
"""

from __future__ import annotations

import sys
from pathlib import Path

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

KCM = Path(__file__).parent.parent
sys.path.insert(0, str(KCM))

from knowable_worlds.dyn_engine import DynSCM                      # noqa: E402
from knowable_worlds.dyn_battery import CHECKPOINTS               # noqa: E402
from knowable_worlds.dyn_confounder import ConfoundedDynSCM       # noqa: E402

INK = RGBColor(0x1B, 0x1B, 0x1F)
MUTED = RGBColor(0x5C, 0x5C, 0x66)
BLUE = RGBColor(0x3A, 0x6E, 0xA5)
GREEN = RGBColor(0x2F, 0x7D, 0x4F)
ORANGE = RGBColor(0xB5, 0x65, 0x2F)
PANEL = RGBColor(0xF4, 0xF4, 0xF2)
FONT = "Segoe UI"

OUT = Path(__file__).parent / "slides"
OUT.mkdir(exist_ok=True)

SW, SH = Inches(13.333), Inches(7.5)


# ---------------------------------------------------------------- figures
def fig_series():
    dyn = DynSCM(n_nodes=8, edge_prob=0.2, seed=300, change_type="sign_flip")
    X = dyn.simulate()
    ce = dyn.changed_edges[0]
    fig, ax = plt.subplots(figsize=(8.6, 3.4), dpi=200)
    t = np.arange(1, dyn.T + 1)
    for k in range(8):
        if k in (ce["i"], ce["j"]):
            continue
        ax.plot(t, X[:, k], color="#d9d9de", lw=0.8, zorder=1)
    ax.plot(t, X[:, ce["i"]], color="#6b6b74", lw=1.4, zorder=2,
            label=f"X{ce['i']+1} (parent of the changed edge)")
    ax.plot(t, X[:, ce["j"]], color="#3a6ea5", lw=2.0, zorder=3,
            label=f"X{ce['j']+1} (affected variable)")
    ax.axvline(dyn.t_change, color="#1b1b1f", ls="--", lw=1.2)
    y_lo, y_hi = ax.get_ylim()          # freeze limits before adding markers
    ax.set_ylim(y_lo, y_hi)
    ax.text(dyn.t_change + 1, y_hi * 0.94,
            f"t* = {dyn.t_change}: edge X{ce['i']+1}→X{ce['j']+1} "
            "reverses sign", fontsize=9)
    ax.plot(CHECKPOINTS, [y_lo] * len(CHECKPOINTS), "^", color="#b5652f",
            markersize=6, clip_on=False, zorder=4, ls="none")
    ax.text(CHECKPOINTS[0] - 2, y_lo, "checkpoints  ", fontsize=8.5,
            color="#b5652f", va="center", ha="right")
    ax.set_xlabel("period")
    ax.legend(loc="upper left", fontsize=9, frameon=False)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    p = OUT / "slide_series.png"
    fig.savefig(p, dpi=200)
    plt.close(fig)
    return p


def fig_confounder():
    scm = ConfoundedDynSCM(seed=300)
    X = scm.simulate()
    ck = 55
    mask = np.array([scm.is_intervened(t) for t in range(1, ck + 1)])
    x1, x2 = X[:ck, 0], X[:ck, 1]

    fig, (ax0, ax1) = plt.subplots(
        1, 2, figsize=(8.6, 3.2), dpi=200, width_ratios=[1, 1.5])

    # left: the DAG
    ax0.set_xlim(0, 10); ax0.set_ylim(0, 8); ax0.axis("off")
    def node(x, y, label, dashed=False):
        c = plt.Circle((x, y), 0.95, fc="white",
                       ec="#9a9aa2" if dashed else "#3a6ea5",
                       ls="--" if dashed else "-", lw=1.6, zorder=3)
        ax0.add_patch(c)
        ax0.text(x, y, label, ha="center", va="center", fontsize=11,
                 color="#5c5c66" if dashed else "#1b1b1f", zorder=4)
    def arrow(x0, y0, x1_, y1_, color):
        ax0.annotate("", xy=(x1_, y1_), xytext=(x0, y0),
                     arrowprops=dict(arrowstyle="-|>", color=color, lw=1.8))
    node(3.5, 6.6, "U", dashed=True)
    node(1.5, 2.2, "X1"); node(5.5, 2.2, "X2"); node(8.8, 5.6, "X3")
    arrow(2.9, 5.8, 1.9, 3.1, "#b5652f")
    arrow(4.1, 5.8, 5.1, 3.1, "#b5652f")
    arrow(8.0, 4.9, 6.3, 2.9, "#2f7d4f")
    ax0.plot([2.6, 4.4], [2.2, 2.2], ls=":", color="#b03030", lw=1)
    ax0.text(3.5, 2.45, "×", ha="center", color="#b03030", fontsize=12)
    ax0.text(3.5, 1.6, "no direct edge", ha="center", fontsize=8,
             color="#b03030")
    ax0.text(3.5, 7.85, "U: latent (never shown)", ha="center", fontsize=8.5,
             color="#5c5c66")

    # right: the decoupling in the shown data
    ax1.scatter(x1[~mask], x2[~mask], s=14, color="#3a6ea5", alpha=0.65,
                label=f"observational rows (n={int((~mask).sum())})")
    ax1.scatter(x1[mask], x2[mask], s=22, color="#b5652f", marker="s",
                label=f"labeled do(X1) rows (n={int(mask.sum())})")
    for m, col in ((~mask, "#3a6ea5"), (mask, "#b5652f")):
        b = np.polyfit(x1[m], x2[m], 1)
        xs = np.array([x1.min(), x1.max()])
        ax1.plot(xs, np.polyval(b, xs), color=col, ls="--", lw=1.3)
    ax1.set_xlabel("X1"); ax1.set_ylabel("X2")
    ax1.legend(fontsize=8.5, frameon=False, loc="best")
    ax1.spines[["top", "right"]].set_visible(False)
    ax1.set_title("periods 1–55 of the presented series", fontsize=9,
                  color="#5c5c66")
    fig.tight_layout()
    p = OUT / "slide_confounder.png"
    fig.savefig(p, dpi=200)
    plt.close(fig)
    return p


# ---------------------------------------------------------------- helpers
def _set(run, size, color=INK, bold=False, italic=False, mono=False):
    run.font.name = "Consolas" if mono else FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def new_slide(prs, title=None, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if title:
        tb = s.shapes.add_textbox(Inches(0.55), Inches(0.32),
                                  SW - Inches(1.1), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        _set(tf.paragraphs[0].add_run(), 27, INK, bold=True)
        tf.paragraphs[0].runs[0].text = title
        if subtitle:
            p = tf.add_paragraph()
            r = p.add_run(); r.text = subtitle
            _set(r, 13, MUTED, italic=True)
        # rule under the title
        ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.22),
                                SW - Inches(1.1), Pt(2.2))
        ln.fill.solid(); ln.fill.fore_color.rgb = BLUE
        ln.line.fill.background()
    return s


def bullets(slide, items, left=0.6, top=1.5, width=12.1, height=5.6,
            size=17, gap=6):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        level, text = (item if isinstance(item, tuple) else (0, item))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        segs = text if isinstance(text, list) else [text]
        for seg in segs:
            t, st = (seg if isinstance(seg, tuple) else (seg, {}))
            r = p.add_run(); r.text = t
            _set(r, st.get("size", size - 2 * level),
                 st.get("color", INK if level == 0 else MUTED),
                 bold=st.get("bold", False), italic=st.get("italic", False),
                 mono=st.get("mono", False))
    return tb


def table(slide, rows, col_widths, left=0.6, top=1.55, row_h=0.42,
          size=13, header=True):
    n_r, n_c = len(rows), len(rows[0])
    shp = slide.shapes.add_table(
        n_r, n_c, Inches(left), Inches(top),
        Inches(sum(col_widths)), Inches(row_h * n_r))
    tbl = shp.table
    tbl.first_row = header
    for c, w in enumerate(col_widths):
        tbl.columns[c].width = Inches(w)
    for ri, row in enumerate(rows):
        for ci, cell in enumerate(row):
            c = tbl.cell(ri, ci)
            c.margin_top = c.margin_bottom = Pt(3)
            tf = c.text_frame; tf.word_wrap = True
            segs = cell if isinstance(cell, list) else [cell]
            for seg in segs:
                t, st = (seg if isinstance(seg, tuple) else (seg, {}))
                r = tf.paragraphs[0].add_run(); r.text = t
                _set(r, st.get("size", size),
                     st.get("color", RGBColor(0xFF, 0xFF, 0xFF)
                            if (header and ri == 0) else INK),
                     bold=st.get("bold", header and ri == 0))
            if header and ri == 0:
                c.fill.solid(); c.fill.fore_color.rgb = INK
            elif ri % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = PANEL
            else:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return shp


def note(slide, text, top=6.85, size=11):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top),
                                  SW - Inches(1.2), Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    _set(r, size, MUTED, italic=True)


# ---------------------------------------------------------------- deck
def build():
    p_series = fig_series()
    p_conf = fig_confounder()

    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    # 1 · title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), SW - Inches(1.8),
                              Inches(2.8))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "Knowable Worlds"
    _set(r, 44, INK, bold=True)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = ("Measuring whether LLM forecasts — and LLM causal beliefs — "
              "track a changing causal environment")
    _set(r, 22, BLUE)
    p = tf.add_paragraph(); p.space_before = Pt(18)
    r = p.add_run()
    r.text = "Motivation and study design · Sean Kelley · July 2026"
    _set(r, 14, MUTED)
    ln = s.shapes.add_shape(1, Inches(0.9), Inches(2.25), Inches(2.6), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()

    # 2 · research question
    s = new_slide(prs, "Research question",
                  "forecasting under structural change requires two distinct "
                  "abilities")
    bullets(s, [
        [("1.  Updating the predictive model when the regime shifts.",
          {"bold": True})],
        [("2.  Relating that update to the change in the underlying causal "
          "structure.", {"bold": True})],
        "",
        [("Are these abilities ", {}), ("coupled", {"bold": True,
                                                    "color": GREEN}),
         (" or ", {}), ("dissociated", {"bold": True, "color": ORANGE}),
         (" in current LLMs?", {})],
        (1, "under coupling, a structural change elicits both forecast "
            "adaptation and a corresponding revision of the stated causal "
            "model"),
        (1, "under dissociation, forecasts adapt through recency alone while "
            "the stated causal model remains unrevised — prediction without "
            "understanding"),
        "",
        "The two are observationally equivalent on standard forecasting "
        "benchmarks. Distinguishing them requires scoring both outputs — "
        "stated probabilities and stated structure — against exact ground "
        "truth, before and after a controlled change.",
    ], top=1.6, size=18)

    # 3 · measurement problem
    s = new_slide(prs, "The measurement problem",
                  "why the question requires simulated environments")
    bullets(s, [
        [("Naturalistic benchmarks score probabilities against realized "
          "binary outcomes.", {"bold": True})],
        (1, "the true probability of a real event is unknowable, so "
            "calibration can be assessed only in aggregate, by binning "
            "across items"),
        (1, "a model that lacked the relevant information is "
            "indistinguishable from one that failed the computation"),
        "",
        [("In a simulated environment the generating equations are known.",
          {"bold": True})],
        (1, "the true probability p* of every item is exactly computable; "
            "each stated probability is scored item by item"),
        (1, "the true causal structure is known at every period; stated "
            "structure is scored per edge"),
        (1, "the timing and content of the change are experimenter-"
            "controlled, so predictive and causal-structural updating are "
            "measured separately against the optimum"),
        "",
        "The environments are synthetic, so training-data contamination is "
        "excluded by construction.",
    ], top=1.6, size=18)

    # 4 · related work
    s = new_slide(prs, "Related work",
                  "closest prior lines, and what remains open")
    table(s, [
        ["prior line", "representative work", "gap relative to this study"],
        ["causal-reasoning benchmarks spanning association, intervention, "
         "and counterfactuals",
         "CLadder (2312.04350); Corr2Cause (2306.05836)",
         "binary accuracy rather than probability calibration; static "
         "problems; no forecasting"],
        ["forecast evaluation against simulated ground truth",
         "ForecastBench-Sim (2606.18686)",
         "identified as future work; no causal layer"],
        ["divergence between stated reasoning and behavior",
         "Turpin et al. (2305.04388)",
         "free-text rationales; no exact truth for either channel; no "
         "adaptation over time"],
        ["in-context learning under regime change",
         "2604.16988",
         "trained transformers on synthetic tasks; no prompted LLM; no "
         "structure elicitation"],
        ["LLM time-series forecasting with contextual information",
         "CiK (2410.18959); From News to Forecast (2409.17515)",
         "accuracy on real series; ground-truth probabilities unknowable; "
         "no causal elicitation"],
    ], [3.7, 3.4, 5.0], top=1.55, row_h=0.88, size=12.5)
    note(s, "To our knowledge, no prior study scores predictive updating and "
            "causal-structural updating against exact truth in the same "
            "environment; the dissociation question has not been posed "
            "quantitatively.")

    # 5 · the world
    s = new_slide(prs, "Design: the data-generating process",
                  "a lag-1 linear system with a controlled structural change")
    s.shapes.add_picture(str(p_series), Inches(0.6), Inches(1.5),
                         width=Inches(8.4))
    bullets(s, [
        [("8 variables", {"bold": True}),
         ("; each period is a linear function of the previous period "
          "(18 cross-lag edges plus autoregressive terms)", {})],
        [("at period t* = 60, ", {}),
         ("one edge changes", {"bold": True, "color": ORANGE}),
         (": addition, removal, sign reversal, or weight doubling", {})],
        [("12 scenarios", {"bold": True}),
         (": 4 change types × 3 independently generated systems", {})],
        [("exact ground truth", {"bold": True}),
         (": the one-step-ahead distribution is Gaussian given the "
          "observed state", {})],
        [("8 checkpoints", {"bold": True}),
         (": 3 pre-change (beliefs must be established before "
          "perseverance is interpretable), 5 post-change", {})],
    ], left=9.15, top=1.55, width=3.7, size=13, gap=10)
    note(s, "Shown: scenario seed 300 (sign reversal). Triangles mark the "
            "elicitation checkpoints.")

    # 6 · elicitation protocol
    s = new_slide(prs, "Elicitation protocol",
                  "the prompt is minimal by design")
    bullets(s, [
        [("The model receives only the raw series through period t, and a "
          "single question.", {"bold": True})],
        (1, "no structural information is provided, and no indication that "
            "a change may occur — unprompted detection is part of the "
            "measurement"),
        (1, "checkpoints are independent calls; no information persists "
            "between them except under explicit carryover conditions"),
    ], top=1.5, height=1.9, size=17)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(3.4), Inches(11.5),
                              Inches(2.9))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate([
        "SYSTEM: You are an expert probability estimator. [...] Respond with "
        "ONLY valid JSON: {\"probability\": <float>}",
        "",
        "USER: Below is a recording of a system of 8 numeric variables, "
        "observed once per period for 66 consecutive periods.",
        "period       X1       X2       X3   ...",
        "     1     0.83     1.94     1.10   ...",
        "   ...      ...      ...      ...",
        "    66     1.02    -0.44     2.31   ...",
        "",
        "What is the probability that X2 in period 67 (the next period) "
        "will exceed -0.52?",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        _set(r, 12, INK, mono=True)
    bg = s.shapes.add_shape(1, Inches(0.75), Inches(3.3), Inches(11.9),
                            Inches(3.05))
    bg.fill.solid(); bg.fill.fore_color.rgb = PANEL
    bg.line.color.rgb = RGBColor(0xC9, 0xC9, 0xC9)
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    note(s, "The threshold τ is placed at the midpoint of the old-regime and "
            "new-regime predictions; the stated probability therefore "
            "necessarily reveals which regime the model treats as current.")

    # 7 · forecast items + references
    s = new_slide(prs, "Forecast items and reference forecasts",
                  "each response is located between exact references")
    bullets(s, [
        [("Two forecast items per checkpoint: the ", {}),
         ("affected variable", {"bold": True, "color": BLUE}),
         (" (a change to edge i→j alters only the conditional distribution "
          "of j) and an ", {}),
         ("unaffected control variable", {"bold": True}),
         (". Error that rises on both indicates nonspecific degradation "
          "rather than detection.", {})],
    ], top=1.5, height=1.1, size=16)
    table(s, [
        ["reference", "data used", "mechanism known", "role"],
        ["truth p*", "current state only", "the current equations",
         "scoring target"],
        ["stale reference", "current state only",
         "the pre-change equations", "perseverance bound"],
        ["rolling window (20 periods)", "last 20 observations",
         "none — refit each checkpoint", "purely recency-based strategy"],
        ["rolling statistician", "last 20 observations",
         "none — asserts edges at |t| > 2",
         "classical statistical reference"],
    ], [2.9, 2.6, 3.2, 3.4], top=2.75, row_h=0.52, size=13)
    note(s, "The model receives strictly more data than every reference. "
            "Each error decomposes as (model − achievable) + (achievable − "
            "truth), and only the first term reflects on the model. "
            "Detection latency is assessed against certified detectability "
            "in the realized sample, not against t* itself.")

    # 8 · structure elicitation
    s = new_slide(prs, "Structure elicitation",
                  "probabilities on targeted edges rather than the full "
                  "graph")
    bullets(s, [
        [("Whole-graph scores are insensitive to the manipulation. ",
          {"bold": True}),
         ("The two regimes share 17 of 18 edges, so any whole-network "
          "measure is dominated by the unchanged edges.", {})],
        "",
        [("Per-edge probability elicitation on tracked pairs.",
          {"bold": True})],
        (1, "at each checkpoint, one pair per query: the changed edge and "
            "three strength-matched comparison edges (two present in both "
            "regimes, one absent in both)"),
        (1, "two probabilities per pair — that the influence is present, "
            "and that it is positive conditional on presence — scored by "
            "Brier score against the exact truth"),
        (1, "the estimand is a difference-in-differences: movement on the "
            "changed edge across t*, relative to the comparison edges"),
        (1, "querying a pair signals its relevance; the comparison edges "
            "receive the identical signal, which the difference removes"),
    ], top=1.55, size=17)

    # 9 · signatures
    s = new_slide(prs, "Predicted signatures by change type",
                  "the behavior of a calibrated respondent on the changed "
                  "edge")
    table(s, [
        ["change at t*", "predicted signature"],
        ["edge added", "P(present) rises from low to high"],
        ["edge removed", "P(present) falls from high to low"],
        ["sign reversal",
         "P(present) remains high; P(positive) crosses 0.5"],
        ["weight doubled",
         "the edge's presence is unchanged, but the evidence for it "
         "strengthens; P(present) should therefore increase toward "
         "certainty"],
    ], [2.6, 9.5], top=1.7, row_h=0.72, size=15)
    note(s, "Weight-doubling scenarios are admitted only when the reference "
            "statistician's confidence (rolling 20-period |t|, recorded "
            "with every response) increases materially across t*; otherwise "
            "detectability, rather than the model, would be under test.",
         top=5.3, size=13)
    bullets(s, [
        [("Forecasts remain the primary channel for magnitude changes; "
          "structure elicitation provides the corresponding belief measure "
          "on the same event.", {"italic": True, "color": MUTED})],
    ], top=6.3, height=0.8, size=14)

    # 10 · multi-edge extension
    s = new_slide(prs, "Extension 1: simultaneous structural changes",
                  "k ∈ {1, 3, 6} edges change at t*, as nested sets")
    bullets(s, [
        [("The manipulation separates two abilities rather than scaling a "
          "single difficulty:", {"bold": True})],
        "",
        [("Detection becomes easier. ", {"bold": True, "color": GREEN}),
         ("Additional changed edges produce additional variables with "
          "elevated forecast residuals.", {})],
        [("Attribution becomes harder. ", {"bold": True, "color": ORANGE}),
         ("Identifying which incoming edges changed grows more demanding "
          "with each affected variable.", {})],
        "",
        "Prediction: the forecast–structure dissociation widens with k, as "
        "the gains from easier detection accrue to forecasting while the "
        "costs of attribution accrue to stated structure.",
        "",
        [("Nesting ensures comparability: the k = 1 change is contained in "
          "the k = 3 set, which is contained in the k = 6 set; the k = 1 "
          "environment is identical to the single-change design.",
          {"size": 14, "color": MUTED})],
    ], top=1.6, size=18)

    # 11 · confounder extension
    s = new_slide(prs, "Extension 2: latent confounding",
                  "an environment in which forecast accuracy alone requires "
                  "causal reasoning")
    s.shapes.add_picture(str(p_conf), Inches(0.55), Inches(1.5),
                         width=Inches(8.2))
    bullets(s, [
        [("In the lagged design, structure is identifiable from observation "
          "alone, so a statistical forecaster is the performance ceiling.",
          {"size": 13, "italic": True})],
        [("This environment removes that ceiling. ", {"bold": True,
                                                      "size": 13}),
         ("A latent variable U drives X1 and X2; approximately 25% of "
          "periods carry labeled interventions on X1, severing U→X1. "
          "The resulting decoupling (right) makes the confounding learnable "
          "from the presented data.", {"size": 13})],
    ], left=8.95, top=1.5, width=3.95, size=13, gap=8)
    table(s, [
        ["the same value a, three questions", "correct answer"],
        ["observe X1 = a next period",
         "use a — an observed X1 is evidence about U"],
        ["set X1 = a next period",
         "disregard a — the intervention severs U→X1"],
        ["set X3 = a (a genuine cause)",
         "use a — excludes the heuristic of discounting all interventions"],
    ], [4.4, 5.6], left=0.6, top=4.95, row_h=0.5, size=13)
    note(s, "A respondent without the causal structure must answer the "
            "observational and interventional questions identically, and is "
            "therefore wrong on one of them by a certified margin: items "
            "are admitted only when the two correct answers differ by at "
            "least 0.15 (realized gaps 0.30–0.55). All ground truth is "
            "closed-form and Monte-Carlo validated.")

    # 12 · controls
    s = new_slide(prs, "Controls",
                  "each removes one alternative explanation")
    table(s, [
        ["control", "alternative explanation removed"],
        ["forecasts on an unaffected variable, every checkpoint",
         "nonspecific post-change degradation read as detection"],
        ["no-change environments (identical design, no change)",
         "drift or degradation that occurs regardless of any change"],
        ["repeated identical prompts; temperature 0",
         "response instability attributable to sampling variance"],
        ["strength-matched comparison edges",
         "global drift in stated beliefs read as edge-specific revision"],
        ["carryover of the correct pre-change structure",
         "failure to revise attributable to the absence of any initial "
         "belief"],
        ["interventions on a genuine cause (set X3)",
         "indiscriminate discounting of interventions scored as causal "
         "competence"],
    ], [5.3, 6.8], top=1.6, row_h=0.62, size=13.5)

    # 13 · hypotheses
    s = new_slide(prs, "Registered hypotheses",
                  "stated prior to data collection")
    table(s, [
        ["", "prediction"],
        ["D1", "after the change, forecast error rises on the affected "
               "variable specifically, not on the controls"],
        ["D2", "responses initially track the pre-change equations and "
               "depart from them as evidence accumulates; the crossing "
               "point defines the detection time"],
        ["D3", "the model underperforms the purely recency-based reference "
               "fitted to the same data"],
        ["D4", "stated structure registers added and removed edges and "
               "sign reversals; weight doubling is registered through "
               "confidence rather than a change in truth value"],
        ["D5", "primary hypothesis: forecasts and stated structure update "
               "at different rates — adaptation without revision would "
               "constitute prediction without understanding"],
    ], [0.7, 11.4], top=1.65, row_h=0.72, size=14)
    note(s, "Analysis: mixed-effects models with scenario random "
            "intercepts. Detection latencies are certified against what the "
            "realized sample supports, not against t*.")

    # 14 · summary
    s = new_slide(prs, "Summary", "what the design provides")
    bullets(s, [
        [("One environment, two scored channels. ", {"bold": True}),
         ("Stated probabilities are scored against exact p*, and stated "
          "structure against the exact graph, before and after a "
          "controlled change.", {})],
        [("The dissociation is measurable. ", {"bold": True}),
         ("Forecast-recovery latency and belief-revision latency are "
          "separate curves on a common clock.", {})],
        [("Failures are attributable. ", {"bold": True}),
         ("Reference forecasts decompose every error into "
          "model-versus-achievable and achievable-versus-truth "
          "components.", {})],
        [("Causal necessity is certified. ", {"bold": True}),
         ("In the confounded environment, the best purely statistical "
          "forecast is wrong by a known per-item margin; only causal "
          "reasoning closes the gap.", {})],
        "",
        [("Status: environments and item batteries implemented and "
          "validated (closed forms verified by Monte Carlo); single-model "
          "pilot complete; multi-model replication, the simultaneous-change "
          "manipulation, and the confounded environment are queued.",
          {"size": 14, "color": MUTED, "italic": True})],
    ], top=1.7, size=18, gap=12)

    out = OUT / "knowable_worlds_motivation_methods.pptx"
    prs.save(out)
    print(f"wrote {out} ({len(prs.slides._sldIdLst)} slides)")
    return out


if __name__ == "__main__":
    build()
